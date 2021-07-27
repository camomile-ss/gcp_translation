#!/usr/bin/python
# coding: utf-8

from pptx import Presentation
import argparse
from gcp_translate_api import gcp_translate_texts

if __name__ == "__main__":

    psr = argparse.ArgumentParser()
    psr.add_argument("infn")
    psr.add_argument("outfn")
    args = psr.parse_args()

    prs = Presentation(args.infn)

    for i, sld in enumerate(prs.slides):  # , start=1):
        print(f"-- {i} --")
        for shp in sld.shapes:
            if shp.has_text_frame:
                texts = shp.text.split("\n")
                texts_jp = gcp_translate_texts(texts)
                shp.text = "\n".join(texts_jp)

    prs.save(args.outfn)
